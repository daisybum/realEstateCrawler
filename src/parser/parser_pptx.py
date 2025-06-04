#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Parser for Real Estate Crawler
Handles extraction of text, tables, and images from PPTX files
Optimized for in-memory processing without saving to disk
"""

import os
import io
import uuid
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple, Union, BinaryIO

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from paddleocr import PaddleOCR
import cv2
import numpy as np
import camelot

from src.config import Config


class PPTXParser:
    """Parser for PPTX files with in-memory processing support"""
    
    def __init__(self, config=None):
        """Initialize the PPTX parser"""
        self.config = config or Config.get_instance()
        self.logger = logging.getLogger(__name__)
        
        # Initialize PaddleOCR for Korean text recognition
        self.ocr = PaddleOCR(use_angle_cls=True, lang='korean', show_log=False)
        
        # Table extraction threshold
        self.table_score_threshold = 0.5  # Minimum accuracy for Camelot tables
        
    def parse_bytes(self, file_obj: BinaryIO) -> Dict[str, Any]:
        """
        Parse PPTX from BytesIO object (alias for extract_content for backward compatibility)
        
        Args:
            file_obj: BytesIO object containing PPTX data
            
        Returns:
            Dictionary with parsed content
        """
        return self.extract_content(file_obj)
        
    def extract_content(self, file_path_or_bytes: Union[str, bytes, BinaryIO]) -> Dict[str, Any]:
        """
        Extract content from a PPTX file or bytes
        
        Args:
            file_path_or_bytes: Path to the PPTX file or bytes/BytesIO object
            
        Returns:
            Dictionary with parsed content
        """
        try:
            # Handle different input types
            if isinstance(file_path_or_bytes, str):
                self.logger.info(f"Parsing PPTX file: {file_path_or_bytes}")
                return self._parse_from_path(file_path_or_bytes)
            elif isinstance(file_path_or_bytes, (bytes, io.BytesIO)):
                self.logger.info("Parsing PPTX from bytes/BytesIO")
                return self._parse_from_bytes(file_path_or_bytes)
            else:
                raise ValueError(f"Unsupported input type: {type(file_path_or_bytes)}")
            
        except Exception as e:
            self.logger.error(f"Error parsing PPTX: {e}")
            return {"error": str(e), "content": f"PPTX 파일 처리 오류: {e}"}
    
    def _parse_from_path(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX from file path"""
        try:
            # Open the file and read as bytes
            with open(file_path, 'rb') as f:
                pptx_bytes = io.BytesIO(f.read())
            
            # Use the bytes parser
            return self._parse_from_bytes(pptx_bytes)
            
        except Exception as e:
            self.logger.error(f"Error reading PPTX file: {e}")
            return {"error": str(e), "content": f"PPTX 파일 읽기 오류: {e}"}
    
    def _parse_from_bytes(self, bytes_data: Union[bytes, BinaryIO]) -> Dict[str, Any]:
        """
        Parse PPTX directly from bytes or BytesIO without saving to disk
        
        Args:
            bytes_data: PPTX data as bytes or BytesIO
            
        Returns:
            Dictionary with parsed content
        """
        result = {
            "content": "",
            "metadata": {},
            "slides": [],
            "elements": []
        }
        
        # Convert to BytesIO if needed
        if isinstance(bytes_data, bytes):
            bytes_io = io.BytesIO(bytes_data)
        else:
            bytes_io = bytes_data
            bytes_io.seek(0)  # Ensure we're at the start of the stream
        
        # Load presentation from BytesIO
        prs = Presentation(bytes_io)
        
        # Extract metadata
        result["metadata"] = self._extract_metadata(prs)
        
        # Extract content from slides
        all_text = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = {
                "slide_num": slide_idx + 1,
                "title": self._get_slide_title(slide),
                "elements": []
            }
            
            for shape_idx, shape in enumerate(slide.shapes):
                element_id = f"{slide_idx}-{shape_idx}"
                
                # Process tables
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_data = self._extract_table(shape)
                    element = {
                        "type": "table",
                        "element_id": element_id,
                        "data": table_data
                    }
                    slide_content["elements"].append(element)
                    result["elements"].append({
                        "slide": slide_idx + 1,
                        **element
                    })
                    
                    # Add table text to content
                    table_text = "\n".join([" | ".join(row) for row in table_data])
                    all_text.append(f"[Table {element_id}]\n{table_text}")
                
                # Process text
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    text_content = shape.text_frame.text.strip()
                    element = {
                        "type": "paragraph",
                        "element_id": element_id,
                        "data": text_content
                    }
                    slide_content["elements"].append(element)
                    result["elements"].append({
                        "slide": slide_idx + 1,
                        **element
                    })
                    
                    # Add text to content
                    all_text.append(text_content)
                
                # Process images
                elif hasattr(shape, "image"):
                    try:
                        # Extract image and run OCR
                        image_data = self._process_image(shape)
                        
                        element = {
                            "type": "image",
                            "element_id": element_id,
                            "ocr_text": image_data["ocr_text"],
                            "tables": image_data["tables"]
                        }
                        slide_content["elements"].append(element)
                        result["elements"].append({
                            "slide": slide_idx + 1,
                            **element
                        })
                        
                        # Add OCR text to content if available
                        if image_data["ocr_text"]:
                            all_text.append(f"[Image {element_id} OCR]\n{image_data['ocr_text']}")
                    
                    except Exception as e:
                        self.logger.warning(f"Error processing image on slide {slide_idx+1}, shape {shape_idx}: {e}")
            
            result["slides"].append(slide_content)
        
        # Combine all text into content
        result["content"] = "\n\n".join(all_text)
        
        return result
            
    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """Extract metadata from a Presentation"""
        metadata = {}
        
        # Extract core properties
        if hasattr(prs, 'core_properties'):
            props = prs.core_properties
            metadata = {
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'keywords': props.keywords or '',
                'created': str(props.created) if props.created else '',
                'modified': str(props.modified) if props.modified else '',
                'last_modified_by': props.last_modified_by or '',
                'slide_count': len(prs.slides)
            }
        else:
            metadata = {
                'slide_count': len(prs.slides)
            }
        
        return metadata
    
    def _get_slide_title(self, slide) -> str:
        """Extract title from a slide"""
        title = ""
        
        # Try to find title in the slide's shapes
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'is_title') and shape.is_title:
                title = shape.text_frame.text
                break
        
        return title.strip()
    
    def _extract_table(self, shape) -> List[List[str]]:
        """Extract table data from a shape"""
        table = shape.table
        return [
            [cell.text_frame.text.strip() for cell in row.cells]
            for row in table.rows
        ]
    
    def _process_image(self, shape) -> Dict[str, Any]:
        """Process an image shape, extract OCR text and detect tables"""
        result = {
            "ocr_text": "",
            "tables": []
        }
        
        try:
            # Convert shape image to PIL Image
            pil_img = self._shape_to_image(shape)
            
            # Run OCR on the image
            cv_img = self._pil_to_cv(pil_img)
            result["ocr_text"] = self._run_paddle_ocr(cv_img)
            
            # Try to extract tables from the image
            result["tables"] = self._extract_tables_from_image(pil_img)
            
        except Exception as e:
            self.logger.warning(f"Error processing image: {e}")
        
        return result
    
    def _shape_to_image(self, shape) -> Image.Image:
        """Convert a shape to PIL Image"""
        return Image.open(io.BytesIO(shape.image.blob))
    
    def _pil_to_cv(self, img: Image.Image) -> np.ndarray:
        """Convert PIL Image to OpenCV format"""
        return cv2.cvtColor(np.asarray(img), cv2.COLOR_RGB2BGR)
    
    def _run_paddle_ocr(self, img: np.ndarray) -> str:
        """
        Run PaddleOCR on an image
        
        Args:
            img: OpenCV image
            
        Returns:
            Extracted text from the image
        """
        try:
            if img is None or img.size == 0:
                return ""
            
            # Check if image is too small for meaningful OCR
            if img.shape[0] < 50 or img.shape[1] < 50:
                return ""  # Skip tiny images like icons
            
            # Preprocess image for better OCR results
            preprocessed_img = self._preprocess_image(img)
            
            # Run PaddleOCR
            ocr_result = self.ocr.ocr(preprocessed_img, cls=True)
            
            if not ocr_result or not ocr_result[0]:
                return ""
            
            # Extract text from OCR result
            text_results = []
            for line in ocr_result[0]:
                if len(line) >= 2 and line[1] and len(line[1]) >= 2:
                    text, confidence = line[1]
                    if confidence > 0.5:  # Only include confident results
                        text_results.append(text)
            
            return " ".join(text_results)
            
        except Exception as e:
            self.logger.error(f"Error running PaddleOCR: {e}")
            return ""
    
    def _preprocess_image(self, img: np.ndarray) -> np.ndarray:
        """
        Preprocess image for better OCR results
        
        Args:
            img: OpenCV image
            
        Returns:
            Preprocessed image
        """
        try:
            # Convert to grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            
            # Apply CLAHE (Contrast Limited Adaptive Histogram Equalization)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            gray = clahe.apply(gray)
            
            # Apply adaptive thresholding
            thresh = cv2.adaptiveThreshold(
                gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                cv2.THRESH_BINARY, 11, 2
            )
            
            # Denoise
            denoised = cv2.fastNlMeansDenoising(thresh, None, 10, 7, 21)
            
            return denoised
            
        except Exception as e:
            self.logger.warning(f"Image preprocessing error: {e}")
            return img
    
    def _extract_tables_from_image(self, img: Image.Image) -> List[List[List[str]]]:
        """
        Extract tables from an image using Camelot
        
        Args:
            img: PIL Image
            
        Returns:
            List of tables, each table is a list of rows
        """
        tables_list = []
        
        try:
            # Convert image to temporary PDF
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                img.convert("RGB").save(tmp_path, "PDF", resolution=300.0)
            
            # Extract tables using Camelot
            camelot_tables = camelot.read_pdf(
                tmp_path, pages="1", flavor="stream"
            )
            
            # Process tables with sufficient accuracy
            for table in camelot_tables:
                if table.accuracy >= self.table_score_threshold * 100:
                    tables_list.append(table.df.values.tolist())
            
            # Clean up temporary file
            os.unlink(tmp_path)
            
        except Exception as e:
            self.logger.warning(f"Error extracting tables from image: {e}")
        
        return tables_list
    
    def _is_image_shape(self, shape) -> bool:
        """
        Check if a shape is an image
        
        Args:
            shape: Shape to check
            
        Returns:
            True if the shape is an image, False otherwise
        """
        try:
            # Check if shape has a picture fill
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            
            # Check if shape has an image property
            if hasattr(shape, 'image'):
                return True
                
            return False
        except Exception:
            return False
    
    def _save_shape_as_image(self, shape, output_path: str) -> bool:
        """
        Save a shape as an image
        
        Args:
            shape: Shape to save
            output_path: Path to save the image to
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # If shape has an image property, save it
            if hasattr(shape, 'image') and shape.image:
                with open(output_path, 'wb') as f:
                    f.write(shape.image.blob)
                return True
                
            return False
        except Exception as e:
            self.logger.warning(f"Error saving shape as image: {e}")
            return False
    
    def _perform_ocr(self, image_path: str) -> str:
        """
        Perform OCR on an image file
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Extracted text from the image
        """
        try:
            # Open the image
            img = Image.open(image_path)
            
            # Convert to OpenCV format for preprocessing
            cv_img = self._pil_to_cv(img)
            
            # Preprocess image for better OCR results
            preprocessed_img = self._preprocess_image(cv_img)
            
            # Perform OCR with PaddleOCR
            return self._run_paddle_ocr(preprocessed_img)
            
        except Exception as e:
            self.logger.warning(f"OCR failed: {e}")
            return ""
    
    def _extract_pptx_content_legacy(self, pptx_path: str) -> Dict[str, Any]:
        """
        Legacy method to extract content from a PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Dictionary containing extracted content
        """
        if not isinstance(pptx_path, str) or not os.path.exists(pptx_path):
            self.logger.error(f"PPTX file not found or invalid path: {pptx_path}")
            return {"error": f"PPTX file not found or invalid path: {pptx_path}"}
        
        try:
            result = {
                "text": "",
                "slides": [],
                "images": [],
                "metadata": {},
                "slide_count": 0
            }
            
            # Extract text, metadata, and images using the new in-memory processing
            return self._parse_from_path(pptx_path)
        
        except Exception as e:
            self.logger.error(f"Error extracting content from PPTX: {e}", exc_info=True)
            return {"error": f"Error extracting content from PPTX: {str(e)}"}
    
    def _extract_pptx_content(self, pptx_path: str) -> Tuple[str, List[Dict[str, Any]], List[Dict[str, Any]], Dict[str, Any], int]:
        """
        Extract content from PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Tuple of (text, slides, images, metadata, slide_count)
        """
        text = ""
        slides = []
        images = []
        metadata = {}
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Open the presentation
            presentation = Presentation(pptx_path)
            
            # Extract metadata
            core_props = presentation.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "last_modified_by": core_props.last_modified_by or "",
                "revision": core_props.revision or 0,
                "category": core_props.category or "",
                "content_status": core_props.content_status or "",
                "language": core_props.language or "",
            }
            
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = ""
                slide_shapes = []
                
                # Process shapes in the slide
                for shape_index, shape in enumerate(slide.shapes):
                    # Extract text from shape
                    if shape.has_text_frame:
                        shape_text = ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                shape_text += run.text + " "
                        
                        slide_text += shape_text.strip() + "\n"
                        
                        slide_shapes.append({
                            "type": "text",
                            "index": shape_index,
                            "text": shape_text.strip()
                        })
                    
                    # Extract tables
                    elif shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = ""
                                for paragraph in cell.text_frame.paragraphs:
                                    cell_text += paragraph.text + " "
                                row_data.append(cell_text.strip())
                            table_data.append(row_data)
                        
                        slide_text += f"[Table with {len(shape.table.rows)} rows and {len(shape.table.columns)} columns]\n"
                        
                        slide_shapes.append({
                            "type": "table",
                            "index": shape_index,
                            "rows": len(shape.table.rows),
                            "columns": len(shape.table.columns),
                            "data": table_data
                        })
                    
                    # Extract images
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Save image to temporary file
                            image_path = os.path.join(temp_dir, f"slide{slide_index+1}_img{shape_index+1}.png")
                            
                            # Get image blob and save it
                            image_blob = shape.image.blob
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_blob)
                            
                            # Perform OCR on the image
                            ocr_text = ""
                            try:
                                # Use our PaddleOCR implementation
                                ocr_text = self._perform_ocr(image_path)
                            except Exception as ocr_e:
                                self.logger.warning(f"OCR failed for image: {ocr_e}")
                            
                            images.append({
                                "slide": slide_index + 1,
                                "index": shape_index + 1,
                                "ocr_text": ocr_text.strip(),
                                "temp_path": image_path
                            })
                            
                            slide_text += f"[Image: {ocr_text.strip() if ocr_text.strip() else 'No text detected'}]\n"
                            
                            slide_shapes.append({
                                "type": "image",
                                "index": shape_index,
                                "ocr_text": ocr_text.strip()
                            })
                        
                        except Exception as img_e:
                            self.logger.warning(f"Could not extract image: {img_e}")
                
                # Add slide information
                slides.append({
                    "index": slide_index + 1,
                    "text": slide_text.strip(),
                    "shapes": slide_shapes
                })
                
                # Add slide text to overall text
                text += f"\n--- Slide {slide_index + 1} ---\n{slide_text}\n"
            
            return text, slides, images, metadata, len(presentation.slides)
        
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content: {e}", exc_info=True)
            return "", [], [], {}, 0
    
    def extract_slide_thumbnails(self, pptx_path: str, output_dir: str, size: Tuple[int, int] = (800, 600)) -> List[str]:
        """
        Extract slide thumbnails from PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            output_dir: Directory to save thumbnails
            size: Thumbnail size as (width, height)
            
        Returns:
            List of paths to the thumbnails
        """
        thumbnail_paths = []
        
        try:
            # Create output directory if it doesn't exist
            os.makedirs(output_dir, exist_ok=True)
            
            # Get the base filename without extension
            base_name = Path(pptx_path).stem
            
            # This is a placeholder since python-pptx doesn't directly support
            # rendering slides as images. In a real implementation, you might:
            # 1. Use a library like comtypes to automate PowerPoint (Windows only)
            # 2. Use a service like LibreOffice to convert PPTX to PDF, then render PDF pages
            # 3. Use a cloud API for conversion
            
            self.logger.warning("Slide thumbnail extraction not fully implemented")
            
            # For now, just return an empty list
            return thumbnail_paths
            
        except Exception as e:
            self.logger.error(f"Error extracting slide thumbnails: {e}", exc_info=True)
            
        return thumbnail_paths
