#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Parser for Real Estate Crawler
Handles extraction of text, tables, and images from PPTX files
"""

import os
import logging
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.config import Config


class PPTXParser:
    """Parser for PPTX files"""
    
    def __init__(self, config=None):
        """Initialize the PPTX parser"""
        self.config = config or Config.get_instance()
        self.logger = logging.getLogger(__name__)
        
    def parse_bytes(self, file_obj) -> Dict[str, Any]:
        """
        Parse PPTX from BytesIO object
        
        Args:
            file_obj: BytesIO object containing PPTX data
            
        Returns:
            Dictionary with parsed content
        """
        try:
            result = {
                "content": "",
                "metadata": {},
                "tables": [],
                "images": []
            }
            
            # Process PPTX
            presentation = Presentation(file_obj)
            
            # Extract text from slides
            text_content = []
            slides_data = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                slide_notes = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract images
                    if self._is_image_shape(shape):
                        try:
                            # Save image to temporary file for OCR
                            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
                                image_path = tmp_file.name
                                self._save_shape_as_image(shape, image_path)
                                
                                # Perform OCR on the image
                                ocr_text = self._perform_ocr(image_path)
                                
                                # Add to images list
                                result["images"].append({
                                    'slide': slide_num + 1,
                                    'index': len(result["images"]),
                                    'ocr_text': ocr_text
                                })
                                
                                # Clean up temporary file
                                os.unlink(image_path)
                        except Exception as e:
                            self.logger.warning(f"Error extracting image from slide {slide_num+1}: {e}")
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide:
                    for shape in slide.notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_notes += shape.text + "\n"
                
                # Add to full text
                slide_content = "\n".join(slide_text)
                text_content.append(f"=== Slide {slide_num+1} ===\n{slide_content}")
                if slide_notes:
                    text_content.append(f"--- Notes ---\n{slide_notes}")
                
                # Add slide data
                slides_data.append({
                    'number': slide_num + 1,
                    'text': slide_content,
                    'notes': slide_notes
                })
            
            result["content"] = "\n\n".join(text_content)
            
            # Extract metadata
            result["metadata"] = {
                'slide_count': len(presentation.slides),
                'slides': slides_data
            }
            
            return result
            
        except Exception as e:
            self.logger.error(f"Error parsing PPTX bytes: {e}")
            return {"error": str(e), "content": f"PPTX 파일 처리 오류: {e}"}
    
    def _is_image_shape(self, shape) -> bool:
        """
        Check if a shape is an image
        
        Args:
            shape: Shape to check
            
        Returns:
            True if the shape is an image, False otherwise
        """
        try:
            # Check if shape has a picture fill
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            
            # Check if shape has an image property
            if hasattr(shape, 'image'):
                return True
                
            return False
        except Exception:
            return False
    
    def _save_shape_as_image(self, shape, output_path: str) -> bool:
        """
        Save a shape as an image
        
        Args:
            shape: Shape to save
            output_path: Path to save the image to
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # If shape has an image property, save it
            if hasattr(shape, 'image') and shape.image:
                with open(output_path, 'wb') as f:
                    f.write(shape.image.blob)
                return True
                
            return False
        except Exception as e:
            self.logger.warning(f"Error saving shape as image: {e}")
            return False
    
    def _perform_ocr(self, image_path: str) -> str:
        """
        Perform OCR on an image file
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Extracted text from the image
        """
        try:
            # Open image with PIL
            img = Image.open(image_path)
            
            # Perform OCR with pytesseract
            text = pytesseract.image_to_string(img, lang='kor+eng')
            return text.strip()
        except Exception as e:
            self.logger.warning(f"OCR failed: {e}")
            return ""
    
    def extract_content(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract content from a PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Dictionary containing extracted content
        """
        if not os.path.exists(pptx_path):
            self.logger.error(f"PPTX file not found: {pptx_path}")
            return {"error": f"PPTX file not found: {pptx_path}"}
        
        try:
            result = {
                "text": "",
                "slides": [],
                "images": [],
                "metadata": {},
                "slide_count": 0
            }
            
            # Extract text, metadata, and images
            text, slides, images, metadata, slide_count = self._extract_pptx_content(pptx_path)
            result["text"] = text
            result["slides"] = slides
            result["images"] = images
            result["metadata"] = metadata
            result["slide_count"] = slide_count
            
            return result
        
        except Exception as e:
            self.logger.error(f"Error extracting content from PPTX: {e}", exc_info=True)
            return {"error": f"Error extracting content from PPTX: {str(e)}"}
    
    def _extract_pptx_content(self, pptx_path: str) -> Tuple[str, List[Dict[str, Any]], List[Dict[str, Any]], Dict[str, Any], int]:
        """
        Extract content from PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Tuple of (text, slides, images, metadata, slide_count)
        """
        text = ""
        slides = []
        images = []
        metadata = {}
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Open the presentation
            presentation = Presentation(pptx_path)
            
            # Extract metadata
            core_props = presentation.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "last_modified_by": core_props.last_modified_by or "",
                "revision": core_props.revision or 0,
                "category": core_props.category or "",
                "content_status": core_props.content_status or "",
                "language": core_props.language or "",
            }
            
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = ""
                slide_shapes = []
                
                # Process shapes in the slide
                for shape_index, shape in enumerate(slide.shapes):
                    # Extract text from shape
                    if shape.has_text_frame:
                        shape_text = ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                shape_text += run.text + " "
                        
                        slide_text += shape_text.strip() + "\n"
                        
                        slide_shapes.append({
                            "type": "text",
                            "index": shape_index,
                            "text": shape_text.strip()
                        })
                    
                    # Extract tables
                    elif shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = ""
                                for paragraph in cell.text_frame.paragraphs:
                                    cell_text += paragraph.text + " "
                                row_data.append(cell_text.strip())
                            table_data.append(row_data)
                        
                        slide_text += f"[Table with {len(shape.table.rows)} rows and {len(shape.table.columns)} columns]\n"
                        
                        slide_shapes.append({
                            "type": "table",
                            "index": shape_index,
                            "rows": len(shape.table.rows),
                            "columns": len(shape.table.columns),
                            "data": table_data
                        })
                    
                    # Extract images
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Save image to temporary file
                            image_path = os.path.join(temp_dir, f"slide{slide_index+1}_img{shape_index+1}.png")
                            
                            # Get image blob and save it
                            image_blob = shape.image.blob
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_blob)
                            
                            # Perform OCR on the image
                            ocr_text = ""
                            try:
                                ocr_text = pytesseract.image_to_string(Image.open(image_path), lang='kor+eng')
                            except Exception as ocr_e:
                                self.logger.warning(f"OCR failed for image: {ocr_e}")
                            
                            images.append({
                                "slide": slide_index + 1,
                                "index": shape_index + 1,
                                "ocr_text": ocr_text.strip(),
                                "temp_path": image_path
                            })
                            
                            slide_text += f"[Image: {ocr_text.strip() if ocr_text.strip() else 'No text detected'}]\n"
                            
                            slide_shapes.append({
                                "type": "image",
                                "index": shape_index,
                                "ocr_text": ocr_text.strip()
                            })
                        
                        except Exception as img_e:
                            self.logger.warning(f"Could not extract image: {img_e}")
                
                # Add slide information
                slides.append({
                    "index": slide_index + 1,
                    "text": slide_text.strip(),
                    "shapes": slide_shapes
                })
                
                # Add slide text to overall text
                text += f"\n--- Slide {slide_index + 1} ---\n{slide_text}\n"
            
            return text, slides, images, metadata, len(presentation.slides)
        
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content: {e}", exc_info=True)
            return "", [], [], {}, 0
    
    def extract_slide_thumbnails(self, pptx_path: str, output_dir: str, size: Tuple[int, int] = (800, 600)) -> List[str]:
        """
        Extract slide thumbnails from PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            output_dir: Directory to save thumbnails
            size: Thumbnail size as (width, height)
            
        Returns:
            List of paths to the thumbnails
        """
        thumbnail_paths = []
        
        try:
            # Create output directory if it doesn't exist
            os.makedirs(output_dir, exist_ok=True)
            
            # Get the base filename without extension
            base_name = Path(pptx_path).stem
            
            # This is a placeholder since python-pptx doesn't directly support
            # rendering slides as images. In a real implementation, you might:
            # 1. Use a library like comtypes to automate PowerPoint (Windows only)
            # 2. Use a service like LibreOffice to convert PPTX to PDF, then render PDF pages
            # 3. Use a cloud API for conversion
            
            self.logger.warning("Slide thumbnail extraction not fully implemented")
            
            # For now, just return an empty list
            return thumbnail_paths
            
        except Exception as e:
            self.logger.error(f"Error extracting slide thumbnails: {e}", exc_info=True)
            
        return thumbnail_paths
