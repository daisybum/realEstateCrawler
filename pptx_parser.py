#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
URL 기반 PPTX → JSONL 변환 스크립트
────────────────────────────────────────────────────────────────
• requests로 PPTX 원본을 메모리(BytesIO)로 로드 → 저장하지 않음
• 이후 로직은 기존과 동일
────────────────────────────────────────────────────────────────
필수 패키지
pip install python-pptx paddleocr paddlepaddle==3.0.0 \
            opencv-python-headless pillow camelot-py[cv] ghostscript \
            pandas tqdm requests
"""

from pathlib import Path
import json, io, tempfile, uuid, requests
from typing import List, Dict, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from paddleocr import PaddleOCR
import cv2
import numpy as np
from PIL import Image

import camelot
from tqdm import tqdm


# ─── 설정 ──────────────────────────────────────────────
PPTX_URL      = "https://cdn.weolbu.com/data_file/202d2a03-248b-499e-867f-4bf7d14cc929.pptx"
OUTPUT_JSONL  = Path("parsed_output.jsonl")      # 결과 JSONL
OCR_LANGS     = "korean"                         # PaddleOCR 언어
TABLE_SCORE_TH = 0.5                             # Camelot 정확도 임계값


# ─── PPTX 원본 다운로드 (메모리) ─────────────────────────
resp = requests.get(PPTX_URL, timeout=30)
resp.raise_for_status()
pptx_stream = io.BytesIO(resp.content)

# ─── 초기화 ────────────────────────────────────────────
ocr = PaddleOCR(lang=OCR_LANGS, show_log=False)   # GPU 사용 시 use_gpu=True


# ─── 도우미 함수들 ────────────────────────────────────
def extract_table(shape) -> List[List[str]]:
    table = shape.table
    return [
        [cell.text_frame.text.strip() for cell in row.cells]
        for row in table.rows
    ]


def shape_to_image(shape) -> Image.Image:
    return Image.open(io.BytesIO(shape.image.blob))


def pil_to_cv(img: Image.Image) -> np.ndarray:
    return cv2.cvtColor(np.asarray(img), cv2.COLOR_RGB2BGR)


def image_to_temp_pdf(img: Image.Image) -> Path:
    tmp = Path(tempfile.gettempdir()) / f"{uuid.uuid4()}.pdf"
    img.convert("RGB").save(tmp, "PDF", resolution=300.0)
    return tmp


# ─── 메인 파싱 루프 ───────────────────────────────────
records: List[Dict[str, Any]] = []
prs = Presentation(pptx_stream)           # 파일이 아닌 BytesIO 객체 사용

for s_idx, slide in enumerate(tqdm(prs.slides, desc="Slides")):
    for shp_idx, shape in enumerate(slide.shapes):
        # 1) 표 Shape ----------------------------------------------------------------
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            records.append({
                "slide": s_idx,
                "element_id": f"{s_idx}-{shp_idx}",
                "type": "table",
                "data": extract_table(shape),
            })
            continue

        # 2) 텍스트 Shape -------------------------------------------------------------
        if shape.has_text_frame and shape.text_frame.text.strip():
            records.append({
                "slide": s_idx,
                "element_id": f"{s_idx}-{shp_idx}",
                "type": "paragraph",
                "data": shape.text_frame.text.strip(),
            })

        # 3) 이미지 Shape -------------------------------------------------------------
        if hasattr(shape, "image"):
            pil_img = shape_to_image(shape)

            # OCR (PIL → OpenCV 배열)
            cv_img = pil_to_cv(pil_img)
            ocr_res = ocr.ocr(cv_img, cls=True)
            ocr_text = (
                " ".join([w[1][0] for w in ocr_res[0]]) if ocr_res else ""
            )

            # Camelot 테이블 검출 (옵션)
            tables_list = []
            try:
                tmp_pdf = image_to_temp_pdf(pil_img)
                camelot_tables = camelot.read_pdf(
                    str(tmp_pdf), pages="1", flavor="stream"
                )
                for tb in camelot_tables:
                    if tb.accuracy >= TABLE_SCORE_TH * 100:
                        tables_list.append(tb.df.values.tolist())
                tmp_pdf.unlink(missing_ok=True)
            except Exception:
                pass

            records.append({
                "slide": s_idx,
                "element_id": f"{s_idx}-{shp_idx}",
                "type": "image",
                "ocr_text": ocr_text,
                "tables": tables_list,
            })


# ─── JSONL 저장 ────────────────────────────────────────
with OUTPUT_JSONL.open("w", encoding="utf-8") as f:
    for rec in records:
        json.dump(rec, f, ensure_ascii=False)
        f.write("\n")

print(f"✅ 파싱 완료 → {OUTPUT_JSONL.resolve()}")
